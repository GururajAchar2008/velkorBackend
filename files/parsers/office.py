"""
files/parsers/office.py
Extract text from DOCX, XLSX and PPTX files.
"""

from typing import Dict, List
import logging
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)

def _ok(text:str, meta=None):
    return {
        "success":True,
        "text":text.strip(),
        "metadata":meta or {},
        "error":None,
    }

def _err(e):
    logger.exception("Office extraction failed")
    return {
        "success":False,
        "text":"",
        "metadata":{},
        "error":str(e),
    }

def extract_docx(file_stream)->Dict:
    try:
        doc=Document(file_stream)
        text="\n".join(p.text for p in doc.paragraphs if p.text)
        meta={"paragraphs":len(doc.paragraphs)}
        return _ok(text,meta)
    except Exception as e:
        return _err(e)

def extract_xlsx(file_stream)->Dict:
    try:
        wb=load_workbook(file_stream,data_only=True)
        lines:List[str]=[]
        for ws in wb.worksheets:
            lines.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals=[str(v) for v in row if v is not None]
                if vals:
                    lines.append(" | ".join(vals))
        return _ok("\n".join(lines),{"sheets":len(wb.sheetnames)})
    except Exception as e:
        return _err(e)

def extract_pptx(file_stream)->Dict:
    try:
        prs=Presentation(file_stream)
        slides=[]
        for i,slide in enumerate(prs.slides,1):
            slides.append(f"# Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape,"text") and shape.text:
                    slides.append(shape.text)
        return _ok("\n".join(slides),{"slides":len(prs.slides)})
    except Exception as e:
        return _err(e)
